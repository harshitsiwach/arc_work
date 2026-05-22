import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_pitch_deck_slides():
    prs = Presentation()
    # Set 16:9 widescreen dimensions (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # BRAND COLORS (Dark Mode Theme)
    # -------------------------------------------------------------
    COLOR_BG = RGBColor(11, 15, 25)         # Slate-950 (Rich Navy-Black)
    COLOR_CARD = RGBColor(22, 28, 45)       # Lighter card container
    COLOR_ACCENT = RGBColor(14, 165, 233)   # Electric Sky Blue
    COLOR_SUCCESS = RGBColor(16, 185, 129)  # Emerald Green
    COLOR_WARNING = RGBColor(239, 68, 68)   # Ruby Red
    COLOR_TEXT_MAIN = RGBColor(255, 255, 255) # Pure White
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # Muted Slate Grey
    COLOR_BORDER = RGBColor(51, 65, 85)     # Slate-700
    COLOR_CARD_EMERALD = RGBColor(16, 45, 35) # Dark emerald for highlights
    COLOR_BORDER_EMERALD = RGBColor(16, 185, 129) # Emerald border

    # -------------------------------------------------------------
    # HELPERS
    # -------------------------------------------------------------
    def apply_bg(slide):
        # Draw background rectangle
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background() # No border
        
        # Glowing top accent line
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = COLOR_SUCCESS
        top_line.line.fill.background()

    def add_slide_header(slide, title, subtitle):
        apply_bg(slide)
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Arial"
        p2.font.size = Pt(15)
        p2.font.color.rgb = COLOR_ACCENT
        p2.space_before = Pt(4)

    def draw_card(slide, x, y, w, h, bg_color=COLOR_CARD, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    def add_card_text(slide, x, y, w, h, title, body, accent_color=COLOR_ACCENT):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.18)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        # Split body by newlines to add clean paragraphs
        body_lines = body.strip().split("\n")
        for line in body_lines:
            p_body = tf.add_paragraph()
            p_body.text = line
            p_body.font.name = "Arial"
            p_body.font.size = Pt(13)
            p_body.font.color.rgb = COLOR_TEXT_MUTED
            p_body.space_before = Pt(6)
            
            # Highlight bullet markers
            if line.startswith("•"):
                p_body.font.size = Pt(12)

    # =============================================================
    # SLIDE 1: COVER & EXECUTIVE SUMMARY
    # =============================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide1)
    
    # Large Cover Title
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = Inches(0)
    
    p = tf1.paragraphs[0]
    p.text = "Arc Work"
    p.font.name = "Arial"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "The On-Chain Operating System for\nInternet Creators & AI Workers"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.space_before = Pt(14)
    
    # Detailed Context on Left
    left_desc = slide1.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.2))
    tf_desc = left_desc.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = Inches(0)
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A decentralized ecosystem designed to eliminate high rent fees, remove settlement holds, and empower AI agents with secure wallet accounts."
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_desc.space_before = Pt(4)
    
    # Right Column: Key Metrics Card (Success highlight)
    draw_card(slide1, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0), bg_color=COLOR_CARD_EMERALD, border_color=COLOR_BORDER_EMERALD)
    
    metric_box = slide1.shapes.add_textbox(Inches(7.2), Inches(1.7), Inches(5.1), Inches(4.6))
    tf_metric = metric_box.text_frame
    tf_metric.word_wrap = True
    tf_metric.margin_left = tf_metric.margin_right = tf_metric.margin_top = tf_metric.margin_bottom = Inches(0.1)
    
    p_m_title = tf_metric.paragraphs[0]
    p_m_title.text = "PLATFORM TRACTION & METRICS"
    p_m_title.font.name = "Arial"
    p_m_title.font.size = Pt(18)
    p_m_title.font.bold = True
    p_m_title.font.color.rgb = COLOR_SUCCESS
    
    metrics = [
        ("2,840+", "Registered Creative Freelancers"),
        ("14.2K+", "Completed Smart Escrow Orders"),
        ("2.5%", "Ultra-Competitive Platform Payout Cut"),
        ("Instant", "Circle USDC/EURC Wallet Settlements")
    ]
    for head, val in metrics:
        p_val = tf_metric.add_paragraph()
        p_val.text = f"{head} — {val}"
        p_val.font.name = "Arial"
        p_val.font.size = Pt(15)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_TEXT_MAIN
        p_val.space_before = Pt(18)

    # =============================================================
    # SLIDE 2: THE PROBLEM
    # =============================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "The Problem", "Centralized bottlenecks restrict freelance gig and creator economic growth")
    
    # 2x2 Grid of Problem Cards
    card_w = Inches(5.6)
    card_h = Inches(2.2)
    gap_x = Inches(0.5)
    gap_y = Inches(0.3)
    
    problems = [
        ("1. Rent Extraction", "• Upwork, Fiverr, and Whop charge 10% to 30% flat fees.\n• Significant FX and cross-border bank charges eat into earnings.\n• High friction margins discourage micro-transactions."),
        ("2. Settlement Hold", "• Web2 payout holds typically take 7-14 business days to clear.\n• Intermediate banks and chargeback vulnerabilities freeze liquidity.\n• Restricts immediate working capital for creative entrepreneurs."),
        ("3. Reputation Lock-In", "• Creator portfolios and ratings are locked inside single platform silos.\n• Zero portability means starting from scratch if switching platforms.\n• Legacy giants capture and own all user reputation metadata."),
        ("4. AI Worker Barriers", "• Standard AI agents lack financial rails to accept and execute contracts.\n• No trust mechanisms to verify AI worker deliverables before payout.\n• Missing link between autonomous workflows and on-chain money.")
    ]
    
    for i, (p_title, p_body) in enumerate(problems):
        row = i // 2
        col = i % 2
        x = Inches(0.8) + col * (card_w + gap_x)
        y = Inches(2.1) + row * (card_h + gap_y)
        
        draw_card(slide2, x, y, card_w, card_h, border_color=COLOR_WARNING)
        add_card_text(slide2, x, y, card_w, card_h, p_title, p_body, accent_color=COLOR_WARNING)

    # =============================================================
    # SLIDE 3: THE SOLUTION
    # =============================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "The Solution", "Low-Fee, Instant on-chain settlements backed by AI verification escrows")
    
    solutions = [
        ("1. Ultra-Low Payouts", "• Flat 2.5% platform transaction fee maximizes earnings.\n• Instant payout settlement in <1 second via Circle USDC/EURC.\n• Eliminates intermediary clearing banks and international wire holds."),
        ("2. Smart Escrows", "• EIP-712 cryptographic templates lock funds securely.\n• Programmable release based on client-defined trigger rules.\n• Reduces chargeback fraud and transaction disputes to zero."),
        ("3. AI Verification Pipeline", "• OpenAI Vision API automatically validates deliverables.\n• Compares visual output directly against contract terms before release.\n• Streamlines quality checks for digital assets and code repositories."),
        ("4. Open Portability", "• Portable identity registry built using ERC-8004 standards.\n• On-chain rating history and reviews follow the user everywhere.\n• Decentralized credentials prevent platform locking and censorship.")
    ]
    
    for i, (s_title, s_body) in enumerate(solutions):
        row = i // 2
        col = i % 2
        x = Inches(0.8) + col * (card_w + gap_x)
        y = Inches(2.1) + row * (card_h + gap_y)
        
        draw_card(slide3, x, y, card_w, card_h, border_color=COLOR_SUCCESS)
        add_card_text(slide3, x, y, card_w, card_h, s_title, s_body, accent_color=COLOR_SUCCESS)

    # =============================================================
    # SLIDE 4: KEY PRODUCT VERTICALS
    # =============================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "Key Product Verticals", "A unified suite of decentralized creative tools")
    
    # 5 Column Grid
    col_w = Inches(2.1)
    col_gap = Inches(0.3)
    col_h = Inches(4.5)
    
    verticals = [
        ("Explore Portal", "`/explore`", "Buy and sell digital templates, graphics, and code scripts. Features seedless checkout powered by Circle developer wallets."),
        ("Freelance Gigs", "`/marketplace`", "Interactive gig board where jobs are locked in smart escrows. AI Vision checks files prior to fund release."),
        ("AI Workers", "`/agents`", "Build, host, and lease autonomous AI agents. Agents receive their own USDC wallets to trade & earn."),
        ("Gated Learning", "`/courses`", "Access premium creative masterclasses using x402 micropayments. Pay-per-module with instant unlocking."),
        ("CCTP Bridge", "`/bridge`", "Circle Cross-Chain Transfer Protocol widget for friction-free stablecoin swaps and bridging across L1/L2s.")
    ]
    
    for i, (v_title, v_route, v_body) in enumerate(verticals):
        x = Inches(0.8) + i * (col_w + col_gap)
        y = Inches(2.1)
        
        draw_card(slide4, x, y, col_w, col_h)
        
        # Custom card content
        box = slide4.shapes.add_textbox(x, y, col_w, col_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.12)
        
        p = tf.paragraphs[0]
        p.text = v_title
        p.font.name = "Arial"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        p_sub = tf.add_paragraph()
        p_sub.text = v_route
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_ACCENT
        p_sub.space_before = Pt(4)
        
        p_body = tf.add_paragraph()
        p_body.text = v_body
        p_body.font.name = "Arial"
        p_body.font.size = Pt(11)
        p_body.font.color.rgb = COLOR_TEXT_MUTED
        p_body.space_before = Pt(10)

    # =============================================================
    # SLIDE 5: TECHNICAL ARCHITECTURE
    # =============================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "Technical Architecture", "Seamless integrations of Web3 wallets, smart contracts, and AI pipelines")
    
    step_w = Inches(2.6)
    step_h = Inches(3.8)
    step_gap = Inches(0.4)
    
    steps = [
        ("Step 1: Auth", "Circle DCW", "Google/Email OAuth provides seedless wallet logins, abstracting Web3 complexities for non-crypto creatives."),
        ("Step 2: Lock", "Circle SCP", "Smart Contract Platform deploys EIP-712 escrows holding client funds in USDC until work is validated."),
        ("Step 3: Audit", "OpenAI Vision", "Automated Node backend executes Visual AI validation checking screenshots and files against terms."),
        ("Step 4: Settle", "x402 Micropay", "Micropayment protocol automatically distributes net funds to creator wallet in <1 second.")
    ]
    
    for i, (s_title, s_tech, s_body) in enumerate(steps):
        x = Inches(0.8) + i * (step_w + step_gap)
        y = Inches(2.3)
        
        draw_card(slide5, x, y, step_w, step_h)
        
        # Text inside step card
        box = slide5.shapes.add_textbox(x, y, step_w, step_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = s_title
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        p_tech = tf.add_paragraph()
        p_tech.text = s_tech
        p_tech.font.name = "Arial"
        p_tech.font.size = Pt(14)
        p_tech.font.bold = True
        p_tech.font.color.rgb = COLOR_ACCENT
        p_tech.space_before = Pt(4)
        
        p_body = tf.add_paragraph()
        p_body.text = s_body
        p_body.font.name = "Arial"
        p_body.font.size = Pt(12)
        p_body.font.color.rgb = COLOR_TEXT_MUTED
        p_body.space_before = Pt(8)
        
        # Draw right arrow between steps (except last step)
        if i < 3:
            arrow_x = x + step_w + Inches(0.05)
            arrow_y = y + Inches(1.8)
            arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.3), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_BORDER
            arrow.line.fill.background()

    # =============================================================
    # SLIDE 6: MARKET OPPORTUNITY & SAVINGS
    # =============================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide6, "Market Opportunity", "High creator fee savings and exponential growth alignment")
    
    comp_w = Inches(3.6)
    comp_h = Inches(4.4)
    comp_gap = Inches(0.4)
    
    # Upwork
    x_up = Inches(0.8)
    draw_card(slide6, x_up, Inches(2.1), comp_w, comp_h, border_color=COLOR_WARNING)
    add_card_text(slide6, x_up, Inches(2.1), comp_w, comp_h, "Upwork (Legacy Freelance)", "• Standard Platform Fee: 20%\n• FX & Wire Fees: 2% to 4%\n• Payout Settling: 5-7 business days\n• Lock-In: High, reputation is non-portable.\n• AI: No automated verification support.", accent_color=COLOR_WARNING)
    
    # Whop
    x_whop = Inches(0.8) + (comp_w + comp_gap)
    draw_card(slide6, x_whop, Inches(2.1), comp_w, comp_h, border_color=COLOR_TEXT_MUTED)
    add_card_text(slide6, x_whop, Inches(2.1), comp_w, comp_h, "Whop (Modern Creator Hub)", "• Standard Platform Fee: 5%\n• FX & Stripe Fees: 2.9% + $0.30\n• Payout Settling: 2-3 business days\n• Lock-In: Medium, limited API portability.\n• AI: No integrated AI work auditing.", accent_color=COLOR_TEXT_MUTED)
    
    # Arc Work
    x_arc = Inches(0.8) + 2 * (comp_w + comp_gap)
    draw_card(slide6, x_arc, Inches(2.1), comp_w, comp_h, bg_color=COLOR_CARD_EMERALD, border_color=COLOR_BORDER_EMERALD)
    add_card_text(slide6, x_arc, Inches(2.1), comp_w, comp_h, "Arc Work (On-chain Platform)", "• Platform Transaction Cut: 2.5%\n• FX & Wire Fees: 0% (USDC Native)\n• Payout Settling: <1 second\n• Lock-In: Zero (ERC-8004 identity)\n• AI: Native Vision API verification.", accent_color=COLOR_SUCCESS)

    # Add Savings Callout Banner
    banner = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_CARD_EMERALD
    banner.line.color.rgb = COLOR_BORDER_EMERALD
    banner.line.width = Pt(1.5)
    
    banner_box = slide6.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.733), Inches(0.8))
    tf_b = banner_box.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.2)
    p_b = tf_b.paragraphs[0]
    p_b.text = "SAVINGS ADVANTAGE: A creator earning $20,000 saves $3,500.00 vs Upwork and $500.00 vs Whop in pure transaction margin."
    p_b.font.name = "Arial"
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_SUCCESS
    p_b.alignment = PP_ALIGN.CENTER

    # =============================================================
    # SLIDE 7: BUSINESS MODEL
    # =============================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide7, "Business Model", "Platform monetization with minimal transaction friction")
    
    models = [
        ("1. Flat Payout Cut", "2.5% Fee Split", "• 1.5% allocated to Net Platform Treasury to fund ecosystem reserves.\n• 1.0% allocated to Relayer gas coverage and network validators.\n• High-margin model driven by automated L2 scaling."),
        ("2. AI Surcharges", "Agent Computing Fee", "• Minor $0.01 - $0.05 surcharge applied per active agent execution step.\n• Offsets OpenAI token costs while retaining clean profit yield.\n• Scalable monetization with autonomous task volumes."),
        ("3. On-chain SaaS", "ERC-8191 Subscriptions", "• Monthly tier options starting at 20 USDC.\n• Grants access to premium seller toolkits, API access, and visual editor templates.\n• Auto-deducted directly from creator's wallet on-chain.")
    ]
    
    for i, (m_title, m_subtitle, m_body) in enumerate(models):
        x = Inches(0.8) + i * (comp_w + comp_gap)
        y = Inches(2.2)
        
        draw_card(slide7, x, y, comp_w, comp_h)
        
        box = slide7.shapes.add_textbox(x, y, comp_w, comp_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.18)
        
        p = tf.paragraphs[0]
        p.text = m_title
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        p_sub = tf.add_paragraph()
        p_sub.text = m_subtitle
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(14)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_ACCENT
        p_sub.space_before = Pt(4)
        
        body_lines = m_body.strip().split("\n")
        for line in body_lines:
            p_body = tf.add_paragraph()
            p_body.text = line
            p_body.font.name = "Arial"
            p_body.font.size = Pt(12)
            p_body.font.color.rgb = COLOR_TEXT_MUTED
            p_body.space_before = Pt(8)

    # =============================================================
    # SLIDE 8: FUTURE ROADMAP
    # =============================================================
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide8, "Future Roadmap", "Direct path to Mainnet launch and fiat integration layers")
    
    roadmap = [
        ("Phase 1: Foundation", "COMPLETED", "• Smart contract escrow release engine on Sepolia/Arc L2.\n• Circle Developer Wallet OAuth client setup.\n• Basic Visual AI validation service endpoint."),
        ("Phase 2: Overhaul", "COMPLETED", "• Collapsible buyer/seller navigation dashboard workspace.\n• OKLCH design system style migration.\n• x402 Learning player framework & tool catalog index."),
        ("Phase 3: Launch", "PLANNED ROADMAP", "• Mainnet deployment & developer SDK release.\n• Circle fiat on/off-ramp native checkout widgets.\n• Cross-platform Zero-Knowledge (ZK) credential reputation export.")
    ]
    
    for i, (r_title, r_status, r_body) in enumerate(roadmap):
        x = Inches(0.8) + i * (comp_w + comp_gap)
        y = Inches(2.2)
        
        # Color coding phase statuses
        b_color = COLOR_SUCCESS if r_status == "COMPLETED" else COLOR_ACCENT
        bg_card_color = COLOR_CARD_EMERALD if r_status == "COMPLETED" else COLOR_CARD
        
        draw_card(slide8, x, y, comp_w, comp_h, bg_color=bg_card_color, border_color=b_color)
        
        box = slide8.shapes.add_textbox(x, y, comp_w, comp_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.18)
        
        p = tf.paragraphs[0]
        p.text = r_title
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        
        p_sub = tf.add_paragraph()
        p_sub.text = r_status
        p_sub.font.name = "Arial"
        p_sub.font.size = Pt(13)
        p_sub.font.bold = True
        p_sub.font.color.rgb = b_color
        p_sub.space_before = Pt(4)
        
        body_lines = r_body.strip().split("\n")
        for line in body_lines:
            p_body = tf.add_paragraph()
            p_body.text = line
            p_body.font.name = "Arial"
            p_body.font.size = Pt(12)
            p_body.font.color.rgb = COLOR_TEXT_MUTED
            p_body.space_before = Pt(8)

    # Save to file
    filepath = "PITCH_DECK.pptx"
    prs.save(filepath)
    print(f"Presentation slides generated successfully and saved to: {filepath}")

if __name__ == "__main__":
    create_pitch_deck_slides()
